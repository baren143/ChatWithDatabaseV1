"""AI-powered report & presentation generation.
User describes what they want in natural language → AI interprets → backend
executes against the FULL dataset → produces a real, data-driven output.
"""

from __future__ import annotations

import io
import json
import logging
from collections import defaultdict
from datetime import datetime
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy import and_, select
from sqlalchemy.orm import Session

from database import get_db
from models import Document, DocumentRow
from dependencies import resolve_user_id_from_request
from routers.filters import fetch_filtered_rows

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api", tags=["reports"])

# ── Request models ─────────────────────────────────────────────────────────────

class ReportRequest(BaseModel):
    prompt: str  # Natural language: "sales by region sorted by revenue"
    document_ids: Optional[List[str]] = None
    output_format: str = "excel"  # excel | csv | pdf


class PresentationRequest(BaseModel):
    prompt: str  # Natural language: "create a presentation about ATM performance"
    document_ids: Optional[List[str]] = None


# ── Shared data helpers ────────────────────────────────────────────────────────

def _get_user_docs(
    db: Session,
    user_id: str,
    document_ids: Optional[List[str]],
) -> List[Document]:
    """Return the user's ready documents."""
    if document_ids:
        docs = list(
            db.execute(
                select(Document).where(
                    and_(
                        Document.user_id == user_id,
                        Document.id.in_(document_ids),
                        Document.status == "ready",
                    )
                )
            ).scalars().all()
        )
    else:
        docs = list(
            db.execute(
                select(Document).where(
                    and_(
                        Document.user_id == user_id,
                        Document.status == "ready",
                    )
                )
            ).scalars().all()
        )
    return docs


def _build_schema_and_stats(db: Session, doc_ids: List[str]) -> Dict[str, Any]:
    """Build schema info + real computed statistics from ALL selected documents."""
    all_rows = list(
        db.execute(
            select(DocumentRow)
            .where(DocumentRow.document_id.in_(doc_ids))
            .order_by(DocumentRow.document_id, DocumentRow.row_index)
        ).scalars().all()
    )
    rows = [r.values for r in all_rows if r.values]
    total_rows = len(rows)

    # Column discovery from sample
    all_columns = set()
    for r in rows[:50]:
        all_columns.update(r.keys())
    columns = list(all_columns)

    # Sample rows (first 5)
    sample_rows = rows[:5]

    # Compute per-column statistics
    col_stats: Dict[str, Dict[str, Any]] = {}
    for col in columns:
        values = [r.get(col) for r in rows if r.get(col) is not None]
        if not values:
            col_stats[col] = {"type": "empty", "count": 0, "unique_count": 0}
            continue

        # Determine if column is numeric
        numeric_vals = []
        for v in values:
            try:
                numeric_vals.append(float(str(v).replace(",", "").strip()))
            except (ValueError, TypeError):
                pass

        is_numeric = len(numeric_vals) > 0.7 * len(values)

        if is_numeric:
            sorted_nums = sorted(numeric_vals)
            col_stats[col] = {
                "type": "numeric",
                "count": len(values),
                "unique_count": len(set(str(v) for v in values)),
                "min": min(sorted_nums),
                "max": max(sorted_nums),
                "sum": sum(sorted_nums),
                "avg": round(sum(sorted_nums) / len(sorted_nums), 2),
                "sample": values[:5],
            }
        else:
            # Categorical — count top values
            val_counts: Dict[str, int] = defaultdict(int)
            for v in values:
                val_counts[str(v).strip()] += 1
            top_values = sorted(val_counts.items(), key=lambda x: -x[1])[:15]
            col_stats[col] = {
                "type": "categorical",
                "count": len(values),
                "unique_count": len(val_counts),
                "top_values": [{"value": v, "count": c} for v, c in top_values],
                "sample": values[:5],
            }

    return {
        "documents": len(doc_ids),
        "total_rows": total_rows,
        "columns": columns,
        "column_stats": col_stats,
        "sample_rows": sample_rows,
    }


def _build_sort_values(rows: List[Dict], sort_col: Optional[str], sort_dir: str) -> List[Dict]:
    """Sort rows by a column, handling numeric vs string sorting."""
    if not sort_col or not rows:
        return rows

    def sort_key(row: Dict) -> Any:
        val = row.get(sort_col)
        if val is None:
            return (0, 0.0, "")
        try:
            num = float(str(val).replace(",", "").strip())
            return (1, num, "")
        except (ValueError, TypeError):
            return (2, 0.0, str(val).lower().strip())

    reverse = sort_dir.lower() == "desc"
    return sorted(rows, key=sort_key, reverse=reverse)


# ── LLM helpers ────────────────────────────────────────────────────────────────

def _call_nvidia_chat(messages: List[Dict], model: str = "meta/llama-3.3-70b-instruct") -> str:
    """Call NVIDIA API for chat completion."""
    import os
    import urllib.request
    import urllib.error

    api_key = os.environ.get("NVIDIA_API_KEY", "")
    if not api_key:
        from config import settings
        api_key = getattr(settings, "nvidia_api_key", "")
    if not api_key:
        raise HTTPException(status_code=500, detail="NVIDIA API key not configured")

    payload = {
        "model": model,
        "messages": messages,
        "temperature": 0.1,
        "max_tokens": 2048,
    }

    req = urllib.request.Request(
        "https://integrate.api.nvidia.com/v1/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = json.loads(resp.read().decode("utf-8"))
            return data["choices"][0]["message"]["content"]
    except urllib.error.HTTPError as e:
        error_body = e.read().decode("utf-8")
        logger.error(f"NVIDIA API error: {e.code} {error_body}")
        raise HTTPException(status_code=502, detail="AI service error. Please try again.")
    except Exception as e:
        logger.error(f"NVIDIA API call failed: {e}")
        raise HTTPException(status_code=502, detail="Failed to reach AI service.")


def _extract_json(text: str) -> Dict[str, Any]:
    """Extract JSON object from LLM response."""
    start = text.find("{")
    end = text.rfind("}") + 1
    if start == -1 or end == 0:
        raise ValueError("No JSON found in response")
    return json.loads(text[start:end])


# ── Report planning ────────────────────────────────────────────────────────────

_REPORT_SYSTEM_PROMPT = """\
You are a data analyst that turns natural language requests into executable report
plans. The user wants a report from their uploaded spreadsheet data.

Your job: output a JSON plan that tells the backend exactly what data to pull,
how to filter, group, sort and summarize it.

Output ONLY a single valid JSON object — no markdown fences, no explanation.

Schema:
{
  "report_title": "<concise title based on the prompt>",
  "filters": [
    {"column": "<exact column name from schema>", "operator": "<eq|ne|contains|gt|lt|gte|lte>", "value": "<filter value>"}
  ],
  "group_by": "<column name or null>",
  "sort_column": "<column name or null>",
  "sort_direction": "<asc|desc>",
  "summary": "<1-3 sentence natural language summary of what this report shows>",
  "notes": "<any additional context or caveats>"
}

Rules:
- Use EXACT column names from the schema (case-sensitive).
- Match filter values CASE to sample data shown.
- If user says "top 10 by X", set sort_column=X, sort_direction=desc.
- If user asks for "total X by Y", set group_by=Y.
- If user says "where Region is North", use filter: column=Region, operator=eq, value=North.
- If user asks for a time period like "Q2 2024" but there's no date column, note that in the summary.
- Do NOT invent column names. Only use columns that appear in the schema.
- Use numeric comparison operators (gt/lt/gte/lte) only on numeric columns.
- If no filters or grouping are specified, leave them as null/empty.
"""


def _plan_report(
    prompt: str,
    stats: Dict[str, Any],
) -> Dict[str, Any]:
    """Ask LLM to plan a report from the user prompt."""
    columns_info = []
    for col, info in stats.get("column_stats", {}).items():
        if info["type"] == "categorical":
            top = ", ".join(f"{v['value']}({v['count']})" for v in info.get("top_values", [])[:5])
            columns_info.append(f"  - {col} (categorical, {info['unique_count']} unique vals): {top}")
        elif info["type"] == "numeric":
            columns_info.append(
                f"  - {col} (numeric, min={info.get('min')}, max={info.get('max')}, avg={info.get('avg')})"
            )
        else:
            columns_info.append(f"  - {col} ({info['type']})")

    sample_rows_text = "\n".join(
        str(r) for r in stats.get("sample_rows", [])[:5]
    )

    msgs = [
        {"role": "system", "content": _REPORT_SYSTEM_PROMPT},
        {
            "role": "user",
            "content": (
                f"Schema & Statistics:\n{chr(10).join(columns_info)}\n\n"
                f"Sample Rows:\n{sample_rows_text}\n\n"
                f"User request: {prompt}\n\n"
                f"Output the JSON plan now."
            ),
        },
    ]

    response_text = _call_nvidia_chat(msgs)
    try:
        return _extract_json(response_text)
    except Exception as e:
        logger.error(f"Failed to parse report plan: {e}, response: {response_text}")
        raise HTTPException(status_code=500, detail="Failed to interpret report request. Please rephrase.")


# ── Presentation planning ──────────────────────────────────────────────────────

_PRESENTATION_SYSTEM_PROMPT = """\
You are a presentation designer. The user wants a PowerPoint from their spreadsheet
data. You will receive real computed statistics (NOT just samples). Use the actual
numbers to build slides — do NOT invent data.

Output ONLY a single valid JSON object — no markdown fences.

JSON schema:
{
  "presentation_title": "<catchy descriptive title>",
  "slides": [
    {
      "slide_title": "<title>",
      "content_type": "<title_only|bullet_points|table|chart>",
      "bullets": ["<insight 1>", "<insight 2>", ...],
      "headers": ["<col1>", "<col2>", ...],
      "rows": [["<val1>", "<val2>", ...], ...],
      "chart_title": "<chart label>",
      "chart_labels": ["<label1>", ...],
      "chart_values": [<num1>, <num2>, ...],
      "narration": "<speaker notes>"
    }
  ],
  "reasoning": "<brief explanation>"
}

CRITICAL RULES:
- Use REAL numbers from the statistics provided. Never invent or guess.
- Use EXACT column names from the schema (case-sensitive).
- MAX 6 slides.
- Slide 1 is ALWAYS a title slide (content_type: "title_only").
- Last slide should be "Key Insights" or "Summary".
- For chart slides, use actual computed figures from the stats (totals, averages, counts).
- For table slides, use the real top/bottom values shown in the data.
- For bullet slides, summarize real patterns visible in the statistics.
- Keep tables to max 8 rows.
- Always provide narration/speaker notes for each slide.
"""


def _plan_presentation(
    prompt: str,
    stats: Dict[str, Any],
) -> Dict[str, Any]:
    """Ask LLM to plan a presentation using real computed statistics."""
    columns_info = []
    for col, info in stats.get("column_stats", {}).items():
        if info["type"] == "categorical":
            top = info.get("top_values", [])
            top_str = ", ".join(f"'{v['value']}': {v['count']}" for v in top[:10])
            columns_info.append(f"  - {col}: categorical, {info['unique_count']} unique values. Distribution: {top_str}")
        elif info["type"] == "numeric":
            columns_info.append(
                f"  - {col}: numeric, total={info.get('sum')}, avg={info.get('avg')}, min={info.get('min')}, max={info.get('max')}"
            )
        else:
            columns_info.append(f"  - {col}: {info['type']}")

    sample_rows_text = "\n".join(
        str(r) for r in stats.get("sample_rows", [])[:5]
    )

    msgs = [
        {"role": "system", "content": _PRESENTATION_SYSTEM_PROMPT},
        {
            "role": "user",
            "content": (
                f"Dataset Overview: {stats.get('total_rows')} rows across {stats.get('documents')} document(s)\n\n"
                f"Column Statistics (computed from FULL dataset):\n{chr(10).join(columns_info)}\n\n"
                f"Sample Rows:\n{sample_rows_text}\n\n"
                f"User request: {prompt}\n\n"
                f"Output the JSON presentation plan now, using the real statistics above."
            ),
        },
    ]

    response_text = _call_nvidia_chat(msgs)
    try:
        return _extract_json(response_text)
    except Exception as e:
        logger.error(f"Failed to parse presentation plan: {e}, response: {response_text}")
        raise HTTPException(status_code=500, detail="Failed to interpret presentation request. Please rephrase.")


# ── Excel report generation ───────────────────────────────────────────────────

def _generate_excel_report(
    plan: Dict[str, Any],
    rows: List[Dict[str, Any]],
) -> bytes:
    """Generate a professional multi-sheet Excel report."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise HTTPException(status_code=501, detail="Excel library not installed")

    wb = openpyxl.Workbook()

    # Colors
    HEADER_FILL = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
    HEADER_FONT = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    DATA_FONT = Font(name="Calibri", size=10, color="1F2937")
    TITLE_FONT = Font(name="Calibri", size=16, bold=True, color="1E3A5F")
    SUMMARY_FONT = Font(name="Calibri", size=11, color="374151")
    thin_border = Border(
        left=Side(style="thin", color="D1D5DB"),
        right=Side(style="thin", color="D1D5DB"),
        top=Side(style="thin", color="D1D5DB"),
        bottom=Side(style="thin", color="D1D5DB"),
    )

    # ── Sheet 1: Summary ──
    ws_summary = wb.active
    ws_summary.title = "Summary"
    ws_summary.merge_cells("A1:F1")
    ws_summary["A1"] = plan.get("report_title", "Data Report")
    ws_summary["A1"].font = TITLE_FONT

    ws_summary["A3"] = "Generated: " + datetime.now().strftime("%Y-%m-%d %H:%M UTC")
    ws_summary["A3"].font = Font(name="Calibri", size=9, color="9CA3AF")

    ws_summary["A5"] = plan.get("summary", "No summary available.")
    ws_summary["A5"].font = SUMMARY_FONT
    ws_summary["A5"].alignment = Alignment(wrap_text=True)
    ws_summary.merge_cells("A5:F7")

    if plan.get("notes"):
        ws_summary["A9"] = "Notes: " + plan["notes"]
        ws_summary["A9"].font = Font(name="Calibri", size=9, italic=True, color="6B7280")

    ws_summary["A11"] = f"Total Rows: {len(rows)}"
    ws_summary["A11"].font = Font(name="Calibri", size=10, bold=True, color="1E3A5F")

    ws_summary.column_dimensions["A"].width = 25
    ws_summary.column_dimensions["B"].width = 20
    ws_summary.column_dimensions["C"].width = 20
    ws_summary.column_dimensions["D"].width = 20
    ws_summary.column_dimensions["E"].width = 20
    ws_summary.column_dimensions["F"].width = 20

    # ── Sheet 2: Data ──
    ws_data = wb.create_sheet("Data")
    if rows:
        headers = list(rows[0].keys())
        for ci, h in enumerate(headers, 1):
            cell = ws_data.cell(row=1, column=ci, value=h)
            cell.font = HEADER_FONT
            cell.fill = HEADER_FILL
            cell.alignment = Alignment(horizontal="center", wrap_text=True)
            cell.border = thin_border

        for ri, row in enumerate(rows[:10000], 2):
            for ci, h in enumerate(headers, 1):
                cell = ws_data.cell(row=ri, column=ci, value=str(row.get(h, "")))
                cell.font = DATA_FONT
                cell.border = thin_border
                if ri % 2 == 0:
                    cell.fill = PatternFill(start_color="F9FAFB", end_color="F9FAFB", fill_type="solid")

        # Auto-width for data sheet
        for ci in range(1, len(headers) + 1):
            max_len = len(str(headers[ci - 1]))
            for ri in range(2, min(len(rows) + 2, 50)):
                val = str(rows[ri - 2].get(headers[ci - 1], ""))
                max_len = max(max_len, min(len(val), 40))
            ws_data.column_dimensions[get_column_letter(ci)].width = max_len + 2

        # Freeze header row
        ws_data.freeze_panes = "A2"

        # Auto-filter
        ws_data.auto_filter.ref = f"A1:{get_column_letter(len(headers))}{len(rows) + 1}"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── PDF report generation ──────────────────────────────────────────────────────

def _generate_pdf_report(
    plan: Dict[str, Any],
    rows: List[Dict[str, Any]],
) -> bytes:
    """Generate a styled PDF report using ReportLab."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate,
            Paragraph,
            Spacer,
            Table,
            TableStyle,
        )
    except ImportError:
        raise HTTPException(status_code=501, detail="PDF library not installed")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        topMargin=0.5 * inch,
        bottomMargin=0.5 * inch,
        leftMargin=0.5 * inch,
        rightMargin=0.5 * inch,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "Title_Custom",
        parent=styles["Title"],
        fontSize=18,
        textColor=colors.HexColor("#1E3A5F"),
        spaceAfter=12,
        alignment=1,  # center
    )
    subtitle_style = ParagraphStyle(
        "Subtitle_Custom",
        parent=styles["Normal"],
        fontSize=8,
        textColor=colors.HexColor("#6B7280"),
        alignment=1,
        spaceAfter=16,
    )
    summary_style = ParagraphStyle(
        "Summary_Custom",
        parent=styles["Normal"],
        fontSize=10,
        textColor=colors.HexColor("#374151"),
        spaceAfter=12,
    )

    elements: list = []

    # Title
    title = plan.get("report_title", "Report")
    elements.append(Paragraph(title, title_style))

    # Metadata line
    meta_text = (
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M UTC')}"
        f" &nbsp;|&nbsp; Rows: {len(rows)}"
    )
    elements.append(Paragraph(meta_text, subtitle_style))

    # Summary
    if plan.get("summary"):
        elements.append(Paragraph(plan["summary"], summary_style))
        elements.append(Spacer(1, 12))

    # Data table
    if rows:
        headers = list(rows[0].keys())
        header_bg = colors.HexColor("#1E3A5F")
        header_fg = colors.white
        row_even_bg = colors.HexColor("#F9FAFB")
        row_odd_bg = colors.white
        border_color = colors.HexColor("#D1D5DB")

        table_data = [[Paragraph(f"<b>{str(h)[:20]}</b>", styles["Normal"]) for h in headers]]

        for i, row in enumerate(rows[:200]):
            table_data.append([
                str(row.get(h, ""))[:25] for h in headers
            ])

        col_widths = [doc.width / len(headers)] * len(headers)
        table = Table(table_data, colWidths=col_widths, repeatRows=1)

        style_commands = [
            ("BACKGROUND", (0, 0), (-1, 0), header_bg),
            ("TEXTCOLOR", (0, 0), (-1, 0), header_fg),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 8),
            ("FONTSIZE", (0, 1), (-1, -1), 7),
            ("ALIGN", (0, 0), (-1, 0), "CENTER"),
            ("GRID", (0, 0), (-1, -1), 0.5, border_color),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ]

        # Alternating row fills
        for i in range(1, len(table_data)):
            bg = row_even_bg if i % 2 == 0 else row_odd_bg
            style_commands.append(("BACKGROUND", (0, i), (-1, i), bg))

        table.setStyle(TableStyle(style_commands))
        elements.append(table)

    doc.build(elements)
    buf.seek(0)
    return buf.getvalue()


# ── CSV report generation ──────────────────────────────────────────────────────

def _generate_csv_report(rows: List[Dict[str, Any]]) -> bytes:
    """Generate CSV from rows."""
    import csv
    buf = io.StringIO()
    if rows:
        headers = list(rows[0].keys())
        writer = csv.DictWriter(buf, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)
    return buf.getvalue().encode("utf-8")


# ── PPTX presentation generation ───────────────────────────────────────────────

def _generate_pptx(plan: Dict[str, Any]) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        raise HTTPException(status_code=500, detail="Presentation library not installed.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x3B, 0x82, 0xF6)
    DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
    TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)

    def _add_title_slide(title: str):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Generated from your data  •  {datetime.now().strftime('%Y-%m-%d')}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
        p2.alignment = PP_ALIGN.CENTER

    def _add_content_slide(slide_data: Dict[str, Any]):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Header bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("slide_title", "Slide")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE

        content_type = slide_data.get("content_type", "bullet_points")

        if content_type == "title_only":
            pass

        elif content_type in ("bullet_points", "text_and_data"):
            bullets = slide_data.get("bullets", [])
            if bullets:
                txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                font_size = Pt(20) if content_type == "bullet_points" else Pt(18)
                for i, bullet in enumerate(bullets):
                    p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                    p2.text = f"•  {bullet}"
                    p2.font.size = font_size
                    p2.font.color.rgb = TEXT_DARK
                    p2.space_before = Pt(8)
                    p2.space_after = Pt(4)

        elif content_type == "table":
            headers = slide_data.get("headers", [])
            rows_data = slide_data.get("rows", [])
            if headers and rows_data:
                table_shape = slide.shapes.add_table(
                    len(rows_data) + 1, len(headers),
                    Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5)
                )
                table = table_shape.table

                for col_idx, h in enumerate(headers):
                    cell = table.cell(0, col_idx)
                    cell.text = str(h)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = BLUE
                    cp = cell.text_frame.paragraphs[0]
                    cp.font.size = Pt(12)
                    cp.font.bold = True
                    cp.font.color.rgb = WHITE
                    cp.alignment = PP_ALIGN.CENTER

                for row_idx, row_vals in enumerate(rows_data[:8]):
                    for col_idx, val in enumerate(row_vals):
                        if col_idx < len(headers):
                            cell = table.cell(row_idx + 1, col_idx)
                            cell.text = str(val)[:50]
                            cell.fill.solid()
                            if row_idx % 2 == 0:
                                cell.fill.fore_color.rgb = WHITE
                            else:
                                cell.fill.fore_color.rgb = LIGHT_GRAY
                            cp = cell.text_frame.paragraphs[0]
                            cp.font.size = Pt(11)
                            cp.font.color.rgb = TEXT_DARK

        elif content_type == "chart":
            chart_labels = slide_data.get("chart_labels", [])
            chart_values = slide_data.get("chart_values", [])
            chart_title = slide_data.get("chart_title", "")
            if chart_labels and chart_values:
                try:
                    from pptx.chart.data import CategoryChartData
                    chart_data = CategoryChartData()
                    chart_data.categories = chart_labels
                    chart_data.add_series(chart_title or "Values", chart_values)
                    chart_shape = slide.shapes.add_chart(
                        None, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.8), chart_data
                    )
                    chart = chart_shape.chart
                    chart.has_legend = True
                    chart.legend.include_in_layout = False
                    if chart.has_title:
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
                except Exception as e:
                    logger.warning(f"Chart failed, falling back to bullets: {e}")
                    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
                    tf2 = txBox2.text_frame
                    tf2.word_wrap = True
                    for i, (label, value) in enumerate(zip(chart_labels, chart_values)):
                        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                        p2.text = f"•  {label}: {value}"
                        p2.font.size = Pt(18)
                        p2.font.color.rgb = TEXT_DARK

        # Narration
        narration = slide_data.get("narration", "")
        if narration:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = narration
            except Exception:
                pass

    presentation_title = plan.get("presentation_title", "Data Presentation")
    slides_plan = plan.get("slides", [])

    _add_title_slide(presentation_title)

    for slide_data in slides_plan:
        _add_content_slide(slide_data)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ── Endpoints ──────────────────────────────────────────────────────────────────

@router.post("/reports/generate")
async def generate_report(
    payload: ReportRequest,
    request: Request,
    db: Session = Depends(get_db),
):
    """
    AI-powered report generation.
    User describes what they want → AI interprets → backend
    executes filters/groups/sorts on full dataset → produces report.
    """
    user_id = resolve_user_id_from_request(request, db)

    if not payload.prompt.strip():
        raise HTTPException(status_code=400, detail="Please describe what report you want.")

    docs = _get_user_docs(db, user_id, payload.document_ids)
    if not docs:
        raise HTTPException(status_code=400, detail="No ready documents found. Please upload and process a file first.")

    doc_ids = [d.id for d in docs]

    # Step 1: Compute statistics from ALL documents
    stats = _build_schema_and_stats(db, doc_ids)

    # Step 2: Ask LLM to plan the report
    plan = _plan_report(payload.prompt, stats)

    # Step 3: Execute the plan — fetch & filter rows
    filters = plan.get("filters", []) or []
    rows = fetch_filtered_rows(db, user_id, doc_ids, filters if filters else None)

    # Step 4: Group if needed
    group_by = plan.get("group_by")
    if group_by and rows:
        from collections import defaultdict
        grouped: Dict[str, List] = defaultdict(list)
        for row in rows:
            key = str(row.get(group_by, "(unknown)"))
            grouped[key].append(row)

        # Convert grouped data into summary rows with counts/sums
        summary_rows = []
        for key, group_rows in grouped.items():
            sr = {group_by: key, "_count": len(group_rows)}
            # Add aggregate of first numeric column found
            for col in stats.get("columns", []):
                if stats.get("column_stats", {}).get(col, {}).get("type") == "numeric":
                    try:
                        total = sum(
                            float(str(r.get(col, 0)).replace(",", "").strip())
                            for r in group_rows
                            if r.get(col) is not None
                        )
                        sr[f"total_{col}"] = round(total, 2)
                    except (ValueError, TypeError, KeyError):
                        pass
            summary_rows.append(sr)
        rows = summary_rows

    # Step 5: Sort
    sort_col = plan.get("sort_column")
    sort_dir = plan.get("sort_direction", "desc")
    if sort_col:
        rows = _build_sort_values(rows, sort_col, sort_dir)

    if not rows:
        rows = [{"note": "No data matched your criteria. Try adjusting your request."}]

    # Step 6: Generate output
    fmt = payload.output_format.lower()
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    title = plan.get("report_title", "Report").replace(" ", "_")[:50]

    if fmt == "csv":
        content = _generate_csv_report(rows)
        return StreamingResponse(
            io.BytesIO(content), media_type="text/csv",
            headers={"Content-Disposition": f'attachment; filename="{title}_{ts}.csv"'},
        )

    elif fmt == "pdf":
        content = _generate_pdf_report(plan, rows)
        return StreamingResponse(
            io.BytesIO(content), media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{title}_{ts}.pdf"'},
        )

    else:  # excel
        content = _generate_excel_report(plan, rows)
        return StreamingResponse(
            io.BytesIO(content),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="{title}_{ts}.xlsx"'},
        )


@router.post("/reports/generate-presentation")
async def generate_presentation(
    payload: PresentationRequest,
    request: Request,
    db: Session = Depends(get_db),
):
    """
    AI-powered presentation generation.
    User describes what they want → AI analyzes real statistics
    from ALL documents → produces data-driven slides.
    """
    user_id = resolve_user_id_from_request(request, db)

    if not payload.prompt.strip():
        raise HTTPException(status_code=400, detail="Please describe the presentation you want.")

    docs = _get_user_docs(db, user_id, payload.document_ids)
    if not docs:
        raise HTTPException(status_code=400, detail="No ready documents found. Please upload and process a file first.")

    doc_ids = [d.id for d in docs]

    # Step 1: Compute statistics from ALL documents
    stats = _build_schema_and_stats(db, doc_ids)

    # Step 2: Ask LLM to plan the presentation using real statistics
    plan = _plan_presentation(payload.prompt, stats)

    # Step 3: Generate the PPTX
    content = _generate_pptx(plan)

    presentation_title = plan.get("presentation_title", "Presentation")
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{presentation_title.replace(' ', '_')[:50]}_{ts}.pptx"

    return StreamingResponse(
        io.BytesIO(content),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )